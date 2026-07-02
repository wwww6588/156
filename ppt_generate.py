#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# ppt_generate.py
# 说明：
# - 依赖：python-pptx, pillow
# - 使用：保存为 ppt_generate.py，创建 assets/ 目录 (可放 logo.png, mech_page1.png, mech_page2.png)
# - 运行：pip install python-pptx pillow    python ppt_generate.py
#
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import os
import sys

# ---------- 配置 ----------
OUTPUT = "GutBarrier_IF-MES_Elasticity_Presentation_v1.pptx"
ASSETS_DIR = "assets"  # 可放 logo.png、mech_page1.png、mech_page2.png 等
SLIDE_WIDTH = Inches(13.333)  # 16:9
SLIDE_HEIGHT = Inches(7.5)
FONT_NAME = "Arial"  # 若你希望使用宋体，可改为 "SimSun"（需在系统已安装）
# --------------------------

def ensure_assets_dir():
    if not os.path.exists(ASSETS_DIR):
        os.makedirs(ASSETS_DIR)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

def set_paragraph_format(text_frame, font_name=FONT_NAME, font_size=18, bold=False, color=(0,0,0)):
    for p in text_frame.paragraphs:
        for run in p.runs:
            try:
                run.font.name = font_name
            except Exception:
                pass
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(*color)

def add_slide(title, bullets, notes="", logo_path=None, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    left = Inches(0.5)
    top = Inches(0.4)
    width = Inches(11.5)
    height = Inches(1.0)

    # 标题框
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title
    set_paragraph_format(tf, font_name=FONT_NAME, font_size=28, bold=True, color=(20,50,120))
    tf.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    # 正文框
    body_top = Inches(1.75)
    body = slide.shapes.add_textbox(left, body_top, width, Inches(4.8))
    b_tf = body.text_frame
    b_tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = b_tf.paragraphs[0]
            p.text = b
        else:
            p = b_tf.add_paragraph()
            p.text = b
        p.level = 0
        try:
            p.font.name = FONT_NAME
        except Exception:
            pass
        p.font.size = Pt(18)

    # 右侧占位图（若提供）
    if image_path:
        img_path = os.path.join(ASSETS_DIR, image_path)
        if os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, Inches(12.0 - 3.8), Inches(1.5), width=Inches(3.2))
            except Exception as e:
                print("插入图片时出错:", img_path, e)

    # logo（右上角）
    if logo_path:
        logo_file = os.path.join(ASSETS_DIR, logo_path)
        if os.path.exists(logo_file):
            try:
                slide.shapes.add_picture(logo_file, Inches(11.6), Inches(0.15), width=Inches(1.2))
            except Exception as e:
                print("插入 logo 时出错:", e)

    # 演讲者备注
    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes

# 幻灯片内容（最终稿）
slides = [
    {
        "title": "肠屏障受损、代谢性内毒素血症与肠道菌群弹性在肥胖代谢并发症中的机制研究",
        "bullets": ["基于 IF + MES 的菌群弹性分型与高弹性菌体内干预验证", "作者 / 单位 / 日期"],
        "notes": "本研究聚焦肠道菌群弹性在肥胖代谢病理中的作用，目标是鉴定高弹性菌并验证其抗肥胖机制。",
    },
    {
        "title": "目录 / 研究框架",
        "bullets": ["1. 背景与假设", "2. 板块一：弹性分型", "3. 板块二：高弹性菌干预",
                    "4. 板块三：机制解析（多组学）", "5. 7 栏机制全景图", "6. 结论与展望"],
        "notes": "快速导览演讲结构，提示听众关注“弹性—体内验证—因果回填”三条主线。"
    },
    {
        "title": "研究背景与总体假设",
        "bullets": ["高脂/低纤饮食→肠黏液/紧密连接受损→LPS 等微生物分子入血",
                    "引发慢性低度炎症→胰岛素抵抗与脂肪堆积",
                    "假设：肠道“菌群弹性”决定群落恢复与代谢稳态"],
        "notes": "解释为何“弹性”比单纯丰度/多样性更能反映群落功能修复能力。"
    },
    {
        "title": "实验动物模型与扰动设计",
        "bullets": ["模型：C57BL/6J-Control、BALB/c-Control、Wistar-Control（大鼠）、C57BL/6J-HFD",
                    "统一施加 AL-FA-RE 间歇禁食扰动，采样时间点：AL / FA / RE"],
        "notes": "多品系设计增强稳健性，IF 模拟稳态—失调—恢复动态。"
    },
    {
        "title": "测序与数据处理流程",
        "bullets": ["粪便采样（AL/FA/RE）→ PacBio 全长 16S（物种水平）",
                    "CLR 中心对数变换消除闭合偏差 → 代入 MES 计算弹性得分 → IQR 划分高/低/无弹性"],
        "notes": "简述 CLR 必要性与 MES 指标构成（幅度、恢复速率、稳态恢复程度）。"
    },
    {
        "title": "弹性评分方法与分类规则",
        "bullets": ["MES 输入：物种丰度随三点时间序列变化（AL→FA→RE）",
                    "评分要素示例：下降幅度、恢复速率、基线恢复程度",
                    "分类：IQR 四分位法划分高/低/无弹性"],
        "notes": "展示典型时间序列示意图（高弹性/低弹性/无弹性）。"
    },
    {
        "title": "弹性分型结果概览（占位）",
        "bullets": ["总体分类分布（高/低/无）— 各模型对比（占位）", "Control vs HFD：高弹性菌数量与比例差异（占位）", "共有核心高弹性菌候选（Top 列表占位）"],
        "notes": "该页待替换真实图表，当前为占位示意。"
    },
    {
        "title": "组间比较：弹性丧失与核心高弹性菌",
        "bullets": ["正常小鼠优势高弹性菌", "肥胖小鼠中特有或优势的高弹性菌", "两组共有核心高弹性菌", "HFD 中弹性丧失的菌株与功能注释（示例：SCFA 生产菌）"],
        "notes": "突出功能门类（如短链脂肪酸生成）在 HFD 中弹性丧失的意义。"
    },
    {
        "title": "核心高弹性菌筛选与体外制备",
        "bullets": ["从板块一筛选跨 Control/HFD 的共有核心高弹性菌","体外厌氧培养，标准化菌液（CFU/mL）、内毒素检测与冻存批次记录"],
        "notes": "强调体外培养纯度与内毒素检测的重要性。"
    },
    {
        "title": "干预动物实验设计与监测",
        "bullets": ["构建 HFD 小鼠并随机分组：PBS、高弹性菌、低弹性菌、无弹性菌","连续 3 周每日灌胃；每 2 天记录体重与摄食量"],
        "notes": "说明随机化与（如有）盲法实施细节与样本量估算。"
    },
    {
        "title": "干预终点与测量项目",
        "bullets": ["糖代谢：GTT、ITT、AUC","脂质蓄积：肝、eWAT、iBAT 称重与 H&E","血清生化：TC、TG、LDL‑C、HDL‑C、ALT、AST",
                    "肠道屏障：结肠病理与紧密连接蛋白（ZO‑1、occludin）","粪便 16S：定植与群落恢复能力验证"],
        "notes": "组合终点用于建立表型、组织学与微生物学的一致证据链。"
    },
    {
        "title": "干预结果示例（占位）",
        "bullets": ["预期：高弹性菌组体重下降、GTT/ITT 改善、脂肪与肝称重下降、屏障修复、粪便菌群恢复（占位示意）"],
        "notes": "此页为占位示例，真数据替换后为主要证据页。"
    },
    {
        "title": "多组学与分子机制解析设计",
        "bullets": ["代谢组：粪便/血浆 LC‑MS（靶向+非靶向）","组织转录组：结肠/eWAT/iBAT → 差异基因与富集分析",
                    "分子验证：qPCR、Western blot（UCP1、PGC1α、ZO‑1、TNF‑α）","代谢物回填：单一/组合代谢物在 HFD 小鼠中回填验证表型复刻"],
        "notes": "说明多组学互补与回填试验在证明因果关系上的关键作用。"
    },
    {
        "title": "关键通路与分子证据（示例）",
        "bullets": ["代谢组：SCFAs ↑、二级胆汁酸改变、色氨酸代谢物调整（示例）","转录组：脂合成↓、产热相关↑、屏障基因↑、NF‑κB 通路抑制","分子验证：qPCR/WB 验证对应分子表达"],
        "notes": "以具体分子示例说明如何从组学到功能验证完成闭环。"
    },
    {
        "title": "7 栏机制全景图（设计说明）",
        "bullets": ["1) 触发因素（高脂、胆汁酸改变、黏液受损、致病菌扩增）",
                    "2) 肠屏障受损（LPS 穿透）",
                    "3) 代谢性内毒素血症（LPS→TLR4→NF‑κB）",
                    "4) 代谢并发症（胰岛素抵抗、脂肪肝）",
                    "5) 测量与研究挑战（LPS 检测局限、多因素干扰）",
                    "6) 机制双向性（菌群紊乱↔炎症↔脂肪组织闭环）",
                    "7) 当前结论（高弹性菌可修复稳态，但临床证据仍需加强）"],
        "notes": "该页为核心视觉摘要，建议用高分辨率矢量海报分两页展示。"
    },
    {
        "title": "主要结论",
        "bullets": ["基于 IF+MES 可客观分型肠道菌株弹性","高弹性菌通过修复群落稳态并分泌有益代谢物修复屏障、抑制炎症、改善糖脂代谢","代谢物回填可建立菌→代谢物→宿主表型的因果链"],
        "notes": "总结创新点并强调潜在临床/治疗价值。"
    },
    {
        "title": "研究局限与未来方向 / 参考与致谢",
        "bullets": ["局限：样本量、模型向人类转化、长期定植与安全性、LPS 检测技术限制","未来方向：扩大人群验证、长期随访、靶细胞/受体机制研究、临床转化探索"],
        "notes": "致谢实验室成员、协作平台与资助。"
    }
]

ensure_assets_dir()

for s in slides:
    add_slide(s["title"], s["bullets"], notes=s.get("notes",""), logo_path="logo.png" if os.path.exists(os.path.join(ASSETS_DIR,"logo.png")) else None, image_path=None)

# 机制海报占位页（分两页）
slide = prs.slides.add_slide(prs.slide_layouts[6])
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(11.5), Inches(1.0))
t.text_frame.text = "机制全景图（1/2） — 触发因素 / 肠屏障 / 内毒素血症 / 并发症"
set_paragraph_format(t.text_frame, font_size=24, bold=True)
if os.path.exists(os.path.join(ASSETS_DIR,"mech_page1.png")):
    slide.shapes.add_picture(os.path.join(ASSETS_DIR,"mech_page1.png"), Inches(0.5), Inches(1.6), width=Inches(12.3))

slide = prs.slides.add_slide(prs.slide_layouts[6])
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(11.5), Inches(1.0))
t.text_frame.text = "机制全景图（2/2） — 测量挑战 / 双向性 / 结论"
set_paragraph_format(t.text_frame, font_size=24, bold=True)
if os.path.exists(os.path.join(ASSETS_DIR,"mech_page2.png")):
    slide.shapes.add_picture(os.path.join(ASSETS_DIR,"mech_page2.png"), Inches(0.5), Inches(1.6), width=Inches(12.3))

# 保存
try:
    prs.save(OUTPUT)
    print("已生成 PPTX：", OUTPUT)
    print("说明：将封面 logo 命名为", os.path.join(ASSETS_DIR,"logo.png"))
    print("如有机制海报，请命名为", os.path.join(ASSETS_DIR,"mech_page1.png / mech_page2.png"))
except Exception as e:
    print("保存 PPTX 时出错:", e)
    sys.exit(1)
